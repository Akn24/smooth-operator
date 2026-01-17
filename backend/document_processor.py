"""Document processing module for extracting text from various file formats."""

import io
import os
import tempfile
import asyncio
from typing import Optional
from dataclasses import dataclass, field
from datetime import datetime
import mimetypes
import logging

logger = logging.getLogger(__name__)


@dataclass
class ExtractedDocument:
    """Represents extracted content from a document."""
    filename: str
    text_content: str
    source_type: str  # 'pdf', 'docx', 'xlsx', 'txt', 'gdoc', 'gsheet', 'gslide'
    extraction_time: datetime = field(default_factory=datetime.utcnow)
    metadata: dict = field(default_factory=dict)
    success: bool = True
    error_message: Optional[str] = None

    @property
    def word_count(self) -> int:
        return len(self.text_content.split()) if self.text_content else 0

    def get_summary(self, max_chars: int = 500) -> str:
        """Get a summary of the document content."""
        if not self.text_content:
            return ""
        return self.text_content[:max_chars] + ("..." if len(self.text_content) > max_chars else "")


class DocumentProcessor:
    """
    Extract text content from various document formats.

    Supports:
    - PDF files (via PyMuPDF/fitz or pdfplumber)
    - Word documents (.docx via python-docx)
    - Excel files (.xlsx via openpyxl)
    - Plain text files
    - Google Docs/Sheets/Slides (via Google APIs)
    """

    SUPPORTED_EXTENSIONS = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'doc',
        '.xlsx': 'xlsx',
        '.xls': 'xls',
        '.txt': 'txt',
        '.md': 'txt',
        '.csv': 'csv',
        '.pptx': 'pptx',
    }

    GOOGLE_MIME_TYPES = {
        'application/vnd.google-apps.document': 'gdoc',
        'application/vnd.google-apps.spreadsheet': 'gsheet',
        'application/vnd.google-apps.presentation': 'gslide',
    }

    EXTRACTION_TIMEOUT = 30  # seconds
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

    def __init__(self, google_credentials=None):
        """
        Initialize the document processor.

        Args:
            google_credentials: Google OAuth credentials for accessing Google Docs/Drive
        """
        self.google_credentials = google_credentials
        self._drive_service = None
        self._docs_service = None
        self._sheets_service = None

    async def extract_from_bytes(
        self,
        content: bytes,
        filename: str,
        mime_type: Optional[str] = None,
    ) -> ExtractedDocument:
        """
        Extract text from file content provided as bytes.

        Args:
            content: File content as bytes
            filename: Original filename (used for extension detection)
            mime_type: Optional MIME type override

        Returns:
            ExtractedDocument with extracted text or error
        """
        # Check file size
        if len(content) > self.MAX_FILE_SIZE:
            return ExtractedDocument(
                filename=filename,
                text_content="",
                source_type="unknown",
                success=False,
                error_message=f"File too large: {len(content)} bytes (max: {self.MAX_FILE_SIZE})",
            )

        # Determine file type
        ext = os.path.splitext(filename.lower())[1]
        source_type = self.SUPPORTED_EXTENSIONS.get(ext, 'unknown')

        if mime_type and mime_type in self.GOOGLE_MIME_TYPES:
            source_type = self.GOOGLE_MIME_TYPES[mime_type]

        try:
            # Use asyncio timeout for extraction
            text = await asyncio.wait_for(
                self._extract_text(content, source_type, filename),
                timeout=self.EXTRACTION_TIMEOUT,
            )

            return ExtractedDocument(
                filename=filename,
                text_content=text,
                source_type=source_type,
                metadata={
                    'file_size': len(content),
                    'mime_type': mime_type,
                },
            )
        except asyncio.TimeoutError:
            return ExtractedDocument(
                filename=filename,
                text_content="",
                source_type=source_type,
                success=False,
                error_message=f"Extraction timeout after {self.EXTRACTION_TIMEOUT}s",
            )
        except Exception as e:
            logger.error(f"Error extracting {filename}: {e}")
            return ExtractedDocument(
                filename=filename,
                text_content="",
                source_type=source_type,
                success=False,
                error_message=str(e),
            )

    async def extract_from_file(self, file_path: str) -> ExtractedDocument:
        """Extract text from a local file."""
        filename = os.path.basename(file_path)
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
            mime_type, _ = mimetypes.guess_type(file_path)
            return await self.extract_from_bytes(content, filename, mime_type)
        except Exception as e:
            return ExtractedDocument(
                filename=filename,
                text_content="",
                source_type="unknown",
                success=False,
                error_message=str(e),
            )

    async def _extract_text(
        self,
        content: bytes,
        source_type: str,
        filename: str,
    ) -> str:
        """Internal method to extract text based on file type."""
        extractors = {
            'pdf': self._extract_pdf,
            'docx': self._extract_docx,
            'doc': self._extract_docx,  # Try docx parser, may not work for old .doc
            'xlsx': self._extract_xlsx,
            'xls': self._extract_xlsx,
            'txt': self._extract_text_file,
            'csv': self._extract_csv,
            'pptx': self._extract_pptx,
        }

        extractor = extractors.get(source_type)
        if extractor:
            return await asyncio.to_thread(extractor, content, filename)

        # Try plain text as fallback
        try:
            return content.decode('utf-8', errors='ignore')
        except Exception:
            return ""

    def _extract_pdf(self, content: bytes, filename: str) -> str:
        """Extract text from PDF using PyMuPDF (fitz)."""
        text_parts = []

        try:
            import fitz  # PyMuPDF

            # Open PDF from bytes
            doc = fitz.open(stream=content, filetype="pdf")

            for page_num, page in enumerate(doc):
                # Extract text with layout preservation
                text = page.get_text("text")
                if text.strip():
                    text_parts.append(f"--- Page {page_num + 1} ---")
                    text_parts.append(text)

                # Also extract tables if present
                tables = page.find_tables()
                if tables:
                    for table in tables:
                        df_text = self._table_to_text(table)
                        if df_text:
                            text_parts.append("\n[Table]")
                            text_parts.append(df_text)

            doc.close()

        except ImportError:
            # Fallback to pdfplumber
            try:
                import pdfplumber

                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            text_parts.append(f"--- Page {page_num + 1} ---")
                            text_parts.append(text)

                        # Extract tables
                        tables = page.extract_tables()
                        for table in tables:
                            if table:
                                text_parts.append("\n[Table]")
                                for row in table:
                                    text_parts.append(" | ".join(str(cell or '') for cell in row))

            except Exception as e:
                logger.error(f"pdfplumber failed: {e}")
                raise

        return "\n".join(text_parts)

    def _table_to_text(self, table) -> str:
        """Convert a fitz table to text."""
        try:
            rows = []
            for row in table.extract():
                if row:
                    rows.append(" | ".join(str(cell or '') for cell in row))
            return "\n".join(rows)
        except Exception:
            return ""

    def _extract_docx(self, content: bytes, filename: str) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            text_parts = []

            for para in doc.paragraphs:
                # Preserve heading structure
                if para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    prefix = '#' * int(level) if level.isdigit() else '#'
                    text_parts.append(f"\n{prefix} {para.text}\n")
                elif para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table in doc.tables:
                text_parts.append("\n[Table]")
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    text_parts.append(" | ".join(cells))

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}")
            raise

    def _extract_xlsx(self, content: bytes, filename: str) -> str:
        """Extract text from Excel file."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

                rows_data = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None for cell in row):
                        row_text = [str(cell) if cell is not None else '' for cell in row]
                        rows_data.append(" | ".join(row_text))

                text_parts.extend(rows_data[:100])  # Limit rows per sheet

                if len(rows_data) > 100:
                    text_parts.append(f"... and {len(rows_data) - 100} more rows")

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error extracting XLSX: {e}")
            raise

    def _extract_text_file(self, content: bytes, filename: str) -> str:
        """Extract text from plain text file."""
        # Try different encodings
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode('utf-8', errors='ignore')

    def _extract_csv(self, content: bytes, filename: str) -> str:
        """Extract text from CSV file."""
        import csv

        text = self._extract_text_file(content, filename)
        text_parts = []

        try:
            reader = csv.reader(io.StringIO(text))
            for i, row in enumerate(reader):
                if i >= 100:  # Limit rows
                    text_parts.append(f"... (truncated, more rows follow)")
                    break
                text_parts.append(" | ".join(row))
        except Exception:
            # Fallback to raw text
            return text[:10000]

        return "\n".join(text_parts)

    def _extract_pptx(self, content: bytes, filename: str) -> str:
        """Extract text from PowerPoint file."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n--- Slide {slide_num} ---\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                    # Handle tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text for cell in row.cells]
                            text_parts.append(" | ".join(cells))

            return "\n".join(text_parts)

        except ImportError:
            logger.warning("python-pptx not installed, skipping PPTX extraction")
            return ""
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            raise

    # ========== Google Docs/Sheets/Slides Extraction ==========

    def _get_drive_service(self):
        """Get or create Google Drive service."""
        if not self._drive_service and self.google_credentials:
            from googleapiclient.discovery import build
            self._drive_service = build('drive', 'v3', credentials=self.google_credentials)
        return self._drive_service

    def _get_docs_service(self):
        """Get or create Google Docs service."""
        if not self._docs_service and self.google_credentials:
            from googleapiclient.discovery import build
            self._docs_service = build('docs', 'v1', credentials=self.google_credentials)
        return self._docs_service

    def _get_sheets_service(self):
        """Get or create Google Sheets service."""
        if not self._sheets_service and self.google_credentials:
            from googleapiclient.discovery import build
            self._sheets_service = build('sheets', 'v4', credentials=self.google_credentials)
        return self._sheets_service

    async def extract_google_doc(self, doc_id: str) -> ExtractedDocument:
        """Extract text from a Google Doc."""
        if not self.google_credentials:
            return ExtractedDocument(
                filename=f"google_doc_{doc_id}",
                text_content="",
                source_type="gdoc",
                success=False,
                error_message="Google credentials not provided",
            )

        try:
            docs_service = self._get_docs_service()
            doc = docs_service.documents().get(documentId=doc_id).execute()

            title = doc.get('title', 'Untitled')
            text_parts = [f"# {title}\n"]

            # Extract content from document body
            content = doc.get('body', {}).get('content', [])
            text_parts.append(self._extract_google_doc_content(content))

            return ExtractedDocument(
                filename=f"{title}.gdoc",
                text_content="\n".join(text_parts),
                source_type="gdoc",
                metadata={'doc_id': doc_id, 'title': title},
            )

        except Exception as e:
            logger.error(f"Error extracting Google Doc {doc_id}: {e}")
            return ExtractedDocument(
                filename=f"google_doc_{doc_id}",
                text_content="",
                source_type="gdoc",
                success=False,
                error_message=str(e),
            )

    def _extract_google_doc_content(self, content: list) -> str:
        """Extract text from Google Doc content elements."""
        text_parts = []

        for element in content:
            if 'paragraph' in element:
                para = element['paragraph']
                para_text = ""

                for elem in para.get('elements', []):
                    if 'textRun' in elem:
                        para_text += elem['textRun'].get('content', '')

                if para_text.strip():
                    # Check for heading style
                    style = para.get('paragraphStyle', {})
                    named_style = style.get('namedStyleType', '')

                    if named_style.startswith('HEADING'):
                        level = named_style.replace('HEADING_', '')
                        prefix = '#' * int(level) if level.isdigit() else '#'
                        text_parts.append(f"{prefix} {para_text.strip()}")
                    else:
                        text_parts.append(para_text)

            elif 'table' in element:
                table = element['table']
                text_parts.append("\n[Table]")
                for row in table.get('tableRows', []):
                    cells = []
                    for cell in row.get('tableCells', []):
                        cell_text = self._extract_google_doc_content(cell.get('content', []))
                        cells.append(cell_text.strip())
                    text_parts.append(" | ".join(cells))

        return "\n".join(text_parts)

    async def extract_google_sheet(self, sheet_id: str) -> ExtractedDocument:
        """Extract text from a Google Sheet."""
        if not self.google_credentials:
            return ExtractedDocument(
                filename=f"google_sheet_{sheet_id}",
                text_content="",
                source_type="gsheet",
                success=False,
                error_message="Google credentials not provided",
            )

        try:
            sheets_service = self._get_sheets_service()

            # Get spreadsheet metadata
            spreadsheet = sheets_service.spreadsheets().get(spreadsheetId=sheet_id).execute()
            title = spreadsheet.get('properties', {}).get('title', 'Untitled')

            text_parts = [f"# {title}\n"]

            # Get data from each sheet
            for sheet in spreadsheet.get('sheets', []):
                sheet_title = sheet.get('properties', {}).get('title', 'Sheet')
                text_parts.append(f"\n=== {sheet_title} ===\n")

                # Get sheet data
                range_name = f"'{sheet_title}'"
                result = sheets_service.spreadsheets().values().get(
                    spreadsheetId=sheet_id,
                    range=range_name,
                ).execute()

                values = result.get('values', [])
                for i, row in enumerate(values[:100]):  # Limit rows
                    text_parts.append(" | ".join(str(cell) for cell in row))

                if len(values) > 100:
                    text_parts.append(f"... and {len(values) - 100} more rows")

            return ExtractedDocument(
                filename=f"{title}.gsheet",
                text_content="\n".join(text_parts),
                source_type="gsheet",
                metadata={'sheet_id': sheet_id, 'title': title},
            )

        except Exception as e:
            logger.error(f"Error extracting Google Sheet {sheet_id}: {e}")
            return ExtractedDocument(
                filename=f"google_sheet_{sheet_id}",
                text_content="",
                source_type="gsheet",
                success=False,
                error_message=str(e),
            )

    async def download_drive_file(self, file_id: str) -> tuple[bytes, str, str]:
        """
        Download a file from Google Drive.

        Returns:
            Tuple of (content_bytes, filename, mime_type)
        """
        if not self.google_credentials:
            raise ValueError("Google credentials not provided")

        drive_service = self._get_drive_service()

        # Get file metadata
        file_metadata = drive_service.files().get(
            fileId=file_id,
            fields='name,mimeType',
        ).execute()

        filename = file_metadata.get('name', 'unknown')
        mime_type = file_metadata.get('mimeType', '')

        # Handle Google Workspace files (export to different format)
        if mime_type in self.GOOGLE_MIME_TYPES:
            export_mime = {
                'application/vnd.google-apps.document': 'text/plain',
                'application/vnd.google-apps.spreadsheet': 'text/csv',
                'application/vnd.google-apps.presentation': 'text/plain',
            }

            content = drive_service.files().export(
                fileId=file_id,
                mimeType=export_mime[mime_type],
            ).execute()

            if isinstance(content, str):
                content = content.encode('utf-8')

            return content, filename, mime_type

        # Download regular file
        from googleapiclient.http import MediaIoBaseDownload

        request = drive_service.files().get_media(fileId=file_id)
        file_buffer = io.BytesIO()
        downloader = MediaIoBaseDownload(file_buffer, request)

        done = False
        while not done:
            _, done = downloader.next_chunk()

        return file_buffer.getvalue(), filename, mime_type


# ========== Utility Functions ==========

def extract_key_metrics(text: str) -> list[str]:
    """
    Extract key numbers and metrics from document text.

    Looks for:
    - Currency amounts ($X, €X)
    - Percentages
    - Large numbers
    - Dates
    """
    import re

    metrics = []

    # Currency patterns
    currency_pattern = r'[\$€£¥]\s*[\d,]+(?:\.\d{2})?(?:\s*(?:million|billion|M|B|K))?'
    currencies = re.findall(currency_pattern, text, re.IGNORECASE)
    metrics.extend(currencies[:10])

    # Percentage patterns
    percent_pattern = r'\d+(?:\.\d+)?\s*%'
    percentages = re.findall(percent_pattern, text)
    metrics.extend(percentages[:10])

    # Large numbers with context
    number_pattern = r'(?:\d{1,3}(?:,\d{3})+|\d+(?:\.\d+)?)\s*(?:million|billion|thousand|M|B|K)'
    numbers = re.findall(number_pattern, text, re.IGNORECASE)
    metrics.extend(numbers[:10])

    return list(set(metrics))


def extract_document_structure(text: str) -> dict:
    """
    Extract structural information from document text.

    Returns:
        Dict with headings, bullet_points, tables count, etc.
    """
    import re

    structure = {
        'headings': [],
        'bullet_points': 0,
        'tables': 0,
        'pages': 0,
    }

    # Find markdown-style headings
    heading_pattern = r'^#+\s+(.+)$'
    headings = re.findall(heading_pattern, text, re.MULTILINE)
    structure['headings'] = headings[:20]

    # Count bullet points
    bullet_pattern = r'^[\s]*[-•*]\s+'
    bullets = re.findall(bullet_pattern, text, re.MULTILINE)
    structure['bullet_points'] = len(bullets)

    # Count tables
    table_pattern = r'\[Table\]'
    tables = re.findall(table_pattern, text)
    structure['tables'] = len(tables)

    # Count pages
    page_pattern = r'---\s*Page\s+\d+\s*---'
    pages = re.findall(page_pattern, text)
    structure['pages'] = len(pages) if pages else 1

    return structure
